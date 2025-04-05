from flask import Flask, request, jsonify, send_from_directory
from pdf2image import convert_from_bytes
import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
import os
import uuid
import io

app = Flask(__name__)

# Directory to save uploaded assets (images)
UPLOAD_FOLDER = os.path.join('static', 'assets')
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# In-memory store for latest assets
latest_image_filenames = []
latest_text_pages = []

# ----------- Helper Functions -----------

def extract_text_from_pdf(pdf_bytes):
    text_content = []
    with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
        for page in doc:
            text = page.get_text()
            if text.strip():
                text_content.append(text.strip())
    return text_content

def extract_text_from_txt(file_bytes):
    return file_bytes.decode('utf-8').splitlines()

def extract_text_from_pptx(file_bytes):
    text_content = []
    ppt = Presentation(io.BytesIO(file_bytes))
    for slide in ppt.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text.strip())
        if slide_text:
            text_content.append("\n".join(slide_text))
    return text_content

# ----------- API Routes -----------

@app.route('/upload_notes', methods=['POST'])
def upload_notes():
    global latest_image_filenames, latest_text_pages

    file = request.files['file']
    filename = file.filename.lower()
    file_bytes = file.read()

    latest_image_filenames = []
    latest_text_pages = []

    if filename.endswith('.pdf'):
        pages = convert_from_bytes(file_bytes)
        latest_text_pages = extract_text_from_pdf(file_bytes)
        for idx, page in enumerate(pages):
            img_filename = f"{uuid.uuid4()}_{idx}.png"
            path = os.path.join(UPLOAD_FOLDER, img_filename)
            page.save(path, 'PNG')
            latest_image_filenames.append(img_filename)

    elif filename.endswith('.txt'):
        latest_text_pages = extract_text_from_txt(file_bytes)

    elif filename.endswith('.pptx'):
        latest_text_pages = extract_text_from_pptx(file_bytes)

    elif filename.endswith(('.png', '.jpg', '.jpeg')):
        img = Image.open(io.BytesIO(file_bytes))
        img_filename = f"{uuid.uuid4()}.png"
        path = os.path.join(UPLOAD_FOLDER, img_filename)
        img.save(path, 'PNG')
        latest_image_filenames.append(img_filename)

    else:
        return jsonify({"error": "Unsupported file format"}), 400

    return jsonify({"status": "Uploaded successfully"})

@app.route('/upload_notes_result', methods=['GET'])
def upload_notes_result():
    global latest_image_filenames, latest_text_pages
    image_urls = [request.host_url.rstrip('/') + '/static/assets/' + fname for fname in latest_image_filenames]
    return jsonify({
        "images": image_urls,
        "texts": latest_text_pages
    })

@app.route('/static/assets/<filename>')
def serve_asset(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@app.route('/test', methods=['GET'])
def test():
    return "✅ Flask backend is running!"

# ----------- Run Local (optional) -----------

if __name__ == "__main__":
    app.run(debug=True)
